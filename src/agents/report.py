"""Report: build the .pptx deck.

Deterministic python-pptx, no LLM. The prose it presents was already written by
the Narrative agent; asking a model to also lay out slides would just add a way
for the deck to come out different each run.

Slides: title, data overview, EDA charts (one per figure), model comparison,
final metrics, recommendations.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from ..state.schema import Artifact, RunState
from ..tools.codegen import run_dir_for
from .common import log, relative_to_run

NODE = "report"

# Matches the chart palette so the deck and the figures read as one thing.
INK_PRIMARY = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, text, left, top, width, height, size=18, bold=False, color=INK_SECONDARY):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return box


def _bullets(slide, items, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for i, item in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = INK_SECONDARY
        run.font.name = "Segoe UI"
        para.space_after = Pt(8)
    return box


def _heading(slide, text):
    _text(slide, text, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          size=28, bold=True, color=INK_PRIMARY)


def _fit_image(slide, image_path: Path, top=None, max_h=None):
    """Centre an image, scaled to fit the content area."""
    from PIL import Image

    if top is None:
        top = Inches(1.4)
    if max_h is None:
        max_h = Inches(5.4)

    with Image.open(image_path) as img:
        ratio = img.height / img.width

    width = Inches(10.5)
    height = Emu(int(width * ratio))
    if height > max_h:
        height = max_h
        width = Emu(int(height / ratio))
    left = Emu(int((SLIDE_W - width) / 2))
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def report_agent(state: RunState) -> RunState:
    """Assemble the deck from what's in state and on disk."""
    run_dir = run_dir_for(state)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Title ---
    slide = _blank(prs)
    _text(slide, state["user_goal"].capitalize(), Inches(0.9), Inches(2.4),
          Inches(11.5), Inches(1.2), size=40, bold=True, color=INK_PRIMARY)
    _text(slide, "Automated analysis by the Autonomous Data Scientist",
          Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6), size=20, color=ACCENT)
    _text(slide, f"Run {state['run_id']}", Inches(0.9), Inches(4.2),
          Inches(11.5), Inches(0.5), size=14)

    # --- Data overview ---
    eda = state.get("eda_findings")
    slide = _blank(prs)
    _heading(slide, "The data")
    overview = []
    if eda:
        overview.append(f"{eda.n_rows:,} rows x {eda.n_cols} columns after cleaning")
        overview.append(f"Predicting: {eda.target}")
    overview.append(f"Task type: {state.get('task_type')}")
    transformations = state.get("transformations", [])
    if transformations:
        overview.append(f"{len(transformations)} cleaning transformations applied:")
        overview += [f"{t.column}: {t.action}" for t in transformations[:6]]
    _bullets(slide, overview, Inches(0.7), Inches(1.5), Inches(12), Inches(5))

    if eda and eda.summary:
        slide = _blank(prs)
        _heading(slide, "What the data shows")
        _text(slide, eda.summary, Inches(0.7), Inches(1.5), Inches(12), Inches(3),
              size=17)
        if eda.anomalies:
            _bullets(slide, [f"Watch out: {a}" for a in eda.anomalies[:4]],
                     Inches(0.7), Inches(4.3), Inches(12), Inches(2.5))

    # --- One slide per chart ---
    for artifact in state.get("artifacts", []):
        if artifact.kind != "plot":
            continue
        path = run_dir / artifact.path
        if path.suffix != ".png" or not path.exists():
            continue  # HTML-only fallback: nothing to embed
        slide = _blank(prs)
        _heading(slide, artifact.label or path.stem)
        _fit_image(slide, path)

    # --- Model comparison ---
    candidates = state.get("candidate_models", [])
    if candidates:
        slide = _blank(prs)
        _heading(slide, "Models considered")
        rows = [
            f"{c.name}: {c.metric} = {c.baseline_score}  ({c.notes})"
            for c in candidates
        ]
        rows.append(f"Selected: {state.get('chosen_model')}")
        tuning = state.get("tuning_results")
        if tuning and tuning.n_trials:
            rows.append(
                f"Tuned over {tuning.n_trials} Optuna trials: "
                f"{tuning.baseline_score} -> {tuning.best_score} "
                f"({tuning.improvement:+.4f})"
            )
        _bullets(slide, rows, Inches(0.7), Inches(1.5), Inches(12), Inches(5))

    # --- Final metrics ---
    evaluation = state.get("eval_metrics")
    if evaluation:
        slide = _blank(prs)
        _heading(slide, "How well it works")
        headline = next(iter(evaluation.metrics.items())) if evaluation.metrics else None
        if headline:
            _text(slide, f"{headline[1]}", Inches(0.7), Inches(1.4),
                  Inches(4), Inches(1.4), size=64, bold=True, color=ACCENT)
            _text(slide, headline[0], Inches(0.75), Inches(2.7), Inches(4),
                  Inches(0.5), size=18)
        _bullets(
            slide,
            [f"{k} = {v}" for k, v in list(evaluation.metrics.items())[1:]]
            + [evaluation.split],
            Inches(5.2), Inches(1.5), Inches(7.4), Inches(2.5),
        )
        _text(slide, evaluation.interpretation, Inches(0.7), Inches(4.2),
              Inches(12), Inches(2.6), size=15)

    # --- Recommendations ---
    slide = _blank(prs)
    _heading(slide, "What we recommend")
    narrative = state.get("narrative", "")
    paragraphs = [p.strip() for p in narrative.split("\n") if p.strip()]
    _text(slide, "\n\n".join(paragraphs[-2:]) if paragraphs else "No narrative available.",
          Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=15)

    path = run_dir / "report.pptx"
    prs.save(path)

    return RunState(
        current_stage=NODE,
        messages=[log(NODE, f"Built report.pptx with {len(prs.slides)} slides.")],
        artifacts=[
            Artifact(kind="report", path=relative_to_run(path, run_dir),
                     label="PowerPoint report")
        ],
    )
